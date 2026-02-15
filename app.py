from flask import Flask, request, render_template, session, jsonify
from groq import Groq
import plotly.graph_objects as go
import sympy as sp
import json
import re
import numpy as np
from PyPDF2 import PdfReader
from docx import Document
import pptx
import random

app = Flask(__name__)
app.secret_key = "some_secret_key_for_session"

MODEL_NAME = "llama-3.1-8b-instant"


def clean_expression(expr):
    """Convert common math notation to Python syntax"""
    expr = str(expr)
    expr = expr.replace('²', '**2')
    expr = expr.replace('³', '**3')
    expr = expr.replace('⁴', '**4')
    expr = expr.replace('⁵', '**5')
    expr = expr.replace('^', '**')
    expr = expr.replace('×', '*')
    expr = expr.replace('÷', '/')
    expr = expr.replace('√', 'sqrt')
    expr = re.sub(r'(\d)([a-zA-Z])', r'\1*\2', expr)
    expr = re.sub(r'(\d)\(', r'\1*(', expr)
    return expr


@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        api_key = request.form.get('api_key')
        if api_key:
            session['api_key'] = api_key
            return render_template('new_math.html')
    return render_template('index.html')
@app.route('/api_guide')
def api_guide():
    return render_template('api_guide.html')

@app.route('/welcome')
def welcome():
    return render_template('index.html')

@app.route('/solve', methods=['POST'])
def solve():
    if 'api_key' not in session:
        return jsonify({'error': 'API key missing'})

    question = request.form.get('question')
    action = request.form.get('action')

    if not question:
        return jsonify({'error': 'Question is required'})

    client = Groq(api_key=session['api_key'])

    # Route to different handlers based on action
    if action == 'solve':
        return handle_solve_with_visualization(question, client)
    elif action == 'explain':
        return handle_explain_with_visualization(question, client)
    elif action == 'alternative_methods':
        return handle_alternative_methods(question, client)
    elif action == 'practice_similar':
        return handle_practice_similar(question, client)
    elif action == 'common_mistakes':
        return handle_common_mistakes(question, client)
    elif action == 'real_world':
        return handle_real_world_application(question, client)
    elif action == 'difficulty_ladder':
        return handle_difficulty_ladder(question, client)
    elif action == 'tutor_mode':
        return handle_tutor_mode(question, client)
    elif action == 'eli5':
        return handle_eli5(question, client)
    elif action == 'concept_map':
        return handle_concept_map(question, client)
    elif action == 'difficulty_rating':
        return handle_difficulty_rating(question, client)
    elif action == 'worksheet':
        return handle_generate_worksheet(question, client)
    else:
        return jsonify({'error': 'Invalid action'})


def handle_solve_with_visualization(question, client):
    """Solve the problem and automatically generate visualization"""
    try:
        solve_prompt = f"Solve this maths question and provide the answer:\n{question}"
        
        completion = client.chat.completions.create(
            model=MODEL_NAME,
            messages=[{"role": "user", "content": solve_prompt}],
            temperature=0
        )
        
        answer = completion.choices[0].message.content.strip()
        viz_analysis = analyze_for_visualization(question, answer, client)
        
        if viz_analysis.get('can_visualize', False):
            graph_html = create_visualization(viz_analysis, question, answer)
            if graph_html:
                return jsonify({
                    'result': answer,
                    'visualization': graph_html,
                    'viz_type': viz_analysis.get('viz_type', 'graph')
                })
        
        return jsonify({'result': answer})
            
    except Exception as e:
        return jsonify({'error': f'Error: {str(e)}'})


def handle_explain_with_visualization(question, client):
    """Explain the problem step-by-step and generate visualization"""
    try:
        explain_prompt = f"Explain step-by-step solution for:\n{question}"
        
        completion = client.chat.completions.create(
            model=MODEL_NAME,
            messages=[{"role": "user", "content": explain_prompt}],
            temperature=0
        )
        
        explanation = completion.choices[0].message.content.strip()
        viz_analysis = analyze_for_visualization(question, explanation, client)
        
        if viz_analysis.get('can_visualize', False):
            graph_html = create_visualization(viz_analysis, question, explanation)
            if graph_html:
                return jsonify({
                    'result': explanation,
                    'visualization': graph_html,
                    'viz_type': viz_analysis.get('viz_type', 'graph')
                })
        
        return jsonify({'result': explanation})
            
    except Exception as e:
        return jsonify({'error': f'Error: {str(e)}'})


def handle_alternative_methods(question, client):
    """Show 2-3 different methods to solve the same problem"""
    try:
        prompt = f"""
        Show 2-3 DIFFERENT methods to solve this problem: {question}
        
        Format your response as:
        
        METHOD 1: [Method Name]
        [Step-by-step solution using this method]
        
        METHOD 2: [Method Name]
        [Step-by-step solution using this method]
        
        METHOD 3: [Method Name] (if applicable)
        [Step-by-step solution using this method]
        
        WHICH TO USE: Brief note on when each method is best
        """
        
        completion = client.chat.completions.create(
            model=MODEL_NAME,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.3
        )
        
        result = completion.choices[0].message.content.strip()
        
        # Try to create visualization for the problem
        viz_analysis = analyze_for_visualization(question, result, client)
        graph_html = None
        if viz_analysis.get('can_visualize', False):
            graph_html = create_visualization(viz_analysis, question, result)
        
        if graph_html:
            return jsonify({
                'result': result,
                'visualization': graph_html,
                'viz_type': viz_analysis.get('viz_type', 'graph')
            })
        
        return jsonify({'result': result})
        
    except Exception as e:
        return jsonify({'error': f'Error: {str(e)}'})


def handle_practice_similar(question, client):
    """Generate 5 similar practice problems"""
    try:
        prompt = f"""
        Based on this problem: {question}
        
        Generate 5 similar practice problems with INCREASING difficulty.
        For each problem, provide:
        - The problem statement
        - The answer (hidden behind a spoiler marker)
        
        Format:
        PROBLEM 1: [Easy - similar to original]
        [problem text]
        Answer: [answer]
        
        PROBLEM 2: [Medium]
        [problem text]
        Answer: [answer]
        
        ... continue for 5 problems total
        """
        
        completion = client.chat.completions.create(
            model=MODEL_NAME,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7
        )
        
        result = completion.choices[0].message.content.strip()
        return jsonify({'result': result})
        
    except Exception as e:
        return jsonify({'error': f'Error: {str(e)}'})


def handle_common_mistakes(question, client):
    """Show common mistakes students make"""
    try:
        prompt = f"""
        For this problem: {question}
        
        List the 3-5 most COMMON MISTAKES students make when solving this.
        For each mistake:
        1. Show the incorrect approach
        2. Explain WHY it's wrong
        3. Show the correct approach
        
        Format:
        MISTAKE 1: [Brief title]
        Incorrect approach: [what students do wrong]
        Why it's wrong: [explanation]
        Correct approach: [right way]
        
        Continue for all common mistakes...
        """
        
        completion = client.chat.completions.create(
            model=MODEL_NAME,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.3
        )
        
        result = completion.choices[0].message.content.strip()
        return jsonify({'result': result})
        
    except Exception as e:
        return jsonify({'error': f'Error: {str(e)}'})


def handle_real_world_application(question, client):
    """Show real-world applications"""
    try:
        prompt = f"""
        For this math problem: {question}
        
        Explain 3 REAL-WORLD applications where this type of math is actually used.
        Include:
        - Specific profession/industry
        - Concrete example scenario
        - Why this math matters there
        
        Make it interesting and relatable to students!
        
        Format:
        APPLICATION 1: [Field/Industry]
        Scenario: [Concrete example]
        Why it matters: [Explanation]
        
        Continue for 3 applications...
        """
        
        completion = client.chat.completions.create(
            model=MODEL_NAME,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.5
        )
        
        result = completion.choices[0].message.content.strip()
        return jsonify({'result': result})
        
    except Exception as e:
        return jsonify({'error': f'Error: {str(e)}'})


def handle_difficulty_ladder(question, client):
    """Show easier and harder versions"""
    try:
        prompt = f"""
        Based on this problem: {question}
        
        Create a difficulty ladder:
        
        EASIER VERSION (Beginner):
        [Simpler problem with same concept]
        Solution: [brief solution]
        
        CURRENT PROBLEM (Intermediate):
        {question}
        
        HARDER VERSION (Advanced):
        [More complex problem with same concept]
        Solution: [brief solution]
        
        EXPERT VERSION (Challenge):
        [Very difficult problem with same concept]
        Solution: [brief solution]
        """
        
        completion = client.chat.completions.create(
            model=MODEL_NAME,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.5
        )
        
        result = completion.choices[0].message.content.strip()
        return jsonify({'result': result})
        
    except Exception as e:
        return jsonify({'error': f'Error: {str(e)}'})


def handle_tutor_mode(question, client):
    """Give hints instead of direct answers"""
    try:
        prompt = f"""
        Act as a patient tutor. For this problem: {question}
        
        DO NOT give the answer directly. Instead provide:
        
        1. WHAT TYPE OF PROBLEM: Identify what math concept this is
        2. KEY CONCEPT: What principle/formula is needed?
        3. FIRST STEP: What should the student do first? (be specific but don't solve)
        4. HINT FOR MIDDLE: What to watch out for in the middle steps
        5. HOW TO CHECK: How can they verify their answer is correct?
        
        Encourage them to try solving it themselves!
        """
        
        completion = client.chat.completions.create(
            model=MODEL_NAME,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.3
        )
        
        result = completion.choices[0].message.content.strip()
        return jsonify({'result': result})
        
    except Exception as e:
        return jsonify({'error': f'Error: {str(e)}'})


def handle_eli5(question, client):
    """Explain like I'm 5"""
    try:
        prompt = f"""
        Explain this math problem to a 5-year-old: {question}
        
        Use:
        - Simple words
        - Fun analogies (cookies, toys, animals, etc.)
        - Short sentences
        - No jargon
        
        Make it fun and easy to understand!
        Then show the solution in simple terms.
        """
        
        completion = client.chat.completions.create(
            model=MODEL_NAME,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.5
        )
        
        result = completion.choices[0].message.content.strip()
        return jsonify({'result': result})
        
    except Exception as e:
        return jsonify({'error': f'Error: {str(e)}'})


def handle_concept_map(question, client):
    """Show concept map and prerequisites"""
    try:
        prompt = f"""
        For this problem: {question}
        
        Create a concept map showing:
        
        PREREQUISITES (What you need to know first):
        - [Concept 1]
        - [Concept 2]
        - [Concept 3]
        
        THIS PROBLEM USES:
        - [Main concept 1]
        - [Main concept 2]
        - [Main concept 3]
        
        NEXT STEPS (What to learn after this):
        - [Advanced topic 1]
        - [Advanced topic 2]
        - [Advanced topic 3]
        
        GRADE LEVEL: [What grade typically learns this]
        """
        
        completion = client.chat.completions.create(
            model=MODEL_NAME,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.3
        )
        
        result = completion.choices[0].message.content.strip()
        return jsonify({'result': result})
        
    except Exception as e:
        return jsonify({'error': f'Error: {str(e)}'})


def handle_difficulty_rating(question, client):
    """Rate problem difficulty"""
    try:
        prompt = f"""
        Analyze this problem: {question}
        
        Provide:
        
        DIFFICULTY RATING: [1-10 scale]
        
        GRADE LEVEL: [What grade]
        
        TIME ESTIMATE: [How long it should take]
        
        COMPLEXITY FACTORS:
        - [What makes it easy/hard]
        - [Skills required]
        - [Common challenges]
        
        RECOMMENDATION: [Who should attempt this problem]
        """
        
        completion = client.chat.completions.create(
            model=MODEL_NAME,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.3
        )
        
        result = completion.choices[0].message.content.strip()
        return jsonify({'result': result})
        
    except Exception as e:
        return jsonify({'error': f'Error: {str(e)}'})


def handle_generate_worksheet(question, client):
    """Generate a practice worksheet"""
    try:
        prompt = f"""
        Based on this problem: {question}
        
        Create a PRACTICE WORKSHEET with 10 problems of varying difficulty.
        
        Format:
        
        === PRACTICE WORKSHEET ===
        Topic: [Identify the topic]
        
        SECTION A: Easy (Problems 1-3)
        1. [problem]
        2. [problem]
        3. [problem]
        
        SECTION B: Medium (Problems 4-7)
        4. [problem]
        5. [problem]
        6. [problem]
        7. [problem]
        
        SECTION C: Hard (Problems 8-10)
        8. [problem]
        9. [problem]
        10. [problem]
        
        === ANSWER KEY ===
        1. [answer]
        2. [answer]
        ... (all answers)
        """
        
        completion = client.chat.completions.create(
            model=MODEL_NAME,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7
        )
        
        result = completion.choices[0].message.content.strip()
        return jsonify({'result': result})
        
    except Exception as e:
        return jsonify({'error': f'Error: {str(e)}'})


def analyze_for_visualization(question, answer, client):
    """Analyze the question and answer to determine what visualization to create"""
    
    question_lower = question.lower()
    
    # FORCE visualization for common patterns
    if re.search(r'x[\^²³⁴\*]|[xy]\s*=|f\(x\)|plot|graph|solve.*x', question_lower):
        expr_match = re.search(r'[xy]\s*=\s*([^,\n]+)', question)
        if expr_match:
            expression = expr_match.group(1).strip()
        else:
            expr_match = re.search(r'([\d\w\+\-\*\/\^\²³⁴\(\)\s]+[xy][\d\w\+\-\*\/\^\²³⁴\(\)\s]*)', question)
            if expr_match:
                expression = expr_match.group(1).strip()
            else:
                expression = question.replace('solve', '').replace('Solve', '').strip()
        
        return {
            'can_visualize': True,
            'viz_type': 'function',
            'expression': expression,
            'data': '',
            'details': 'function plot'
        }
    
    # Check for shapes
    if any(word in question_lower for word in ['circle', 'square', 'triangle', 'rectangle', 'area', 'perimeter']):
        return {
            'can_visualize': True,
            'viz_type': 'geometric',
            'expression': question,
            'data': '',
            'details': question
        }
    
    # Check for lists of numbers
    numbers = re.findall(r'\d+', question + answer[:200])
    if len(numbers) >= 3:
        if 'mean' in question_lower or 'average' in question_lower or 'distribution' in question_lower:
            return {
                'can_visualize': True,
                'viz_type': 'statistical',
                'expression': '',
                'data': ','.join(numbers),
                'details': 'statistical analysis'
            }
        else:
            return {
                'can_visualize': True,
                'viz_type': 'data',
                'expression': '',
                'data': ','.join(numbers),
                'details': 'data visualization'
            }
    
    return {'can_visualize': False, 'viz_type': 'none'}


def create_visualization(viz_analysis, question, answer):
    """Create visualization based on analysis"""
    
    viz_type = viz_analysis.get('viz_type', 'none')
    
    try:
        if viz_type == 'function':
            return create_function_plot(viz_analysis)
        elif viz_type == 'geometric':
            return create_geometric_plot(viz_analysis)
        elif viz_type == 'data':
            return create_data_plot(viz_analysis)
        elif viz_type == 'statistical':
            return create_statistical_plot(viz_analysis)
        elif viz_type == 'vector':
            return create_vector_plot(viz_analysis)
    except Exception as e:
        print(f"Visualization creation error: {str(e)}")
    
    return None


def create_function_plot(viz_analysis):
    """Create a plot for mathematical functions"""
    try:
        expression = viz_analysis.get('expression', '')
        if not expression:
            return None
            
        expression = clean_expression(expression)
        x = sp.Symbol('x')
        expr = sp.sympify(expression)
        
        x_vals = np.linspace(-10, 10, 400)
        y_vals = []
        
        for val in x_vals:
            try:
                y = float(expr.subs(x, val))
                if abs(y) < 1e6:
                    y_vals.append(y)
                else:
                    y_vals.append(None)
            except:
                y_vals.append(None)
        
        fig = go.Figure()
        fig.add_trace(go.Scatter(
            x=x_vals,
            y=y_vals,
            mode='lines',
            name=f'f(x) = {expression}',
            line=dict(color='#2E86AB', width=3)
        ))
        
        fig.add_hline(y=0, line_dash="dash", line_color="gray", opacity=0.5)
        fig.add_vline(x=0, line_dash="dash", line_color="gray", opacity=0.5)
        
        fig.update_layout(
            title=f'Graph of {expression}',
            xaxis_title='x',
            yaxis_title='f(x)',
            hovermode='x unified',
            template='plotly_white',
            height=500,
            showlegend=True
        )
        
        return fig.to_html(full_html=False)
        
    except Exception as e:
        print(f"Function plot error: {str(e)}")
        return None


def create_geometric_plot(viz_analysis):
    """Create geometric visualizations"""
    try:
        details = viz_analysis.get('details', '').lower()
        expression = viz_analysis.get('expression', '').lower()
        combined = details + " " + expression
        
        fig = go.Figure()
        
        if 'circle' in combined:
            radius = 5
            numbers = re.findall(r'\d+\.?\d*', expression + details)
            if numbers:
                radius = float(numbers[0])
            
            theta = np.linspace(0, 2*np.pi, 100)
            x = radius * np.cos(theta)
            y = radius * np.sin(theta)
            
            fig.add_trace(go.Scatter(
                x=x, y=y,
                mode='lines',
                fill='toself',
                name=f'Circle (r={radius})',
                line=dict(color='#A23B72', width=3),
                fillcolor='rgba(162, 59, 114, 0.3)'
            ))
            
            fig.add_trace(go.Scatter(
                x=[0], y=[0],
                mode='markers',
                marker=dict(size=10, color='#A23B72'),
                name='Center'
            ))
            
        elif 'square' in combined:
            side = 4
            numbers = re.findall(r'\d+\.?\d*', expression + details)
            if numbers:
                side = float(numbers[0])
            
            half = side / 2
            x = [-half, half, half, -half, -half]
            y = [-half, -half, half, half, -half]
            
            fig.add_trace(go.Scatter(
                x=x, y=y,
                mode='lines',
                fill='toself',
                name=f'Square (side={side})',
                line=dict(color='#6A994E', width=3),
                fillcolor='rgba(106, 153, 78, 0.3)'
            ))
            
        elif 'triangle' in combined:
            side = 2
            numbers = re.findall(r'\d+\.?\d*', expression + details)
            if numbers:
                side = float(numbers[0])
            
            height = side * np.sqrt(3) / 2
            x = [-side/2, side/2, 0, -side/2]
            y = [0, 0, height, 0]
            
            fig.add_trace(go.Scatter(
                x=x, y=y,
                mode='lines',
                fill='toself',
                name=f'Triangle (side={side})',
                line=dict(color='#F18F01', width=3),
                fillcolor='rgba(241, 143, 1, 0.3)'
            ))
        
        elif 'rectangle' in combined:
            numbers = re.findall(r'\d+\.?\d*', expression + details)
            width = float(numbers[0]) if len(numbers) > 0 else 6
            height = float(numbers[1]) if len(numbers) > 1 else 4
            
            x = [0, width, width, 0, 0]
            y = [0, 0, height, height, 0]
            
            fig.add_trace(go.Scatter(
                x=x, y=y,
                mode='lines',
                fill='toself',
                name=f'Rectangle ({width}×{height})',
                line=dict(color='#6A994E', width=3),
                fillcolor='rgba(106, 153, 78, 0.3)'
            ))
        
        else:
            return None
        
        fig.update_layout(
            title='Geometric Visualization',
            xaxis_title='x',
            yaxis_title='y',
            template='plotly_white',
            height=500,
            yaxis=dict(scaleanchor="x", scaleratio=1),
            showlegend=True
        )
        
        return fig.to_html(full_html=False)
        
    except Exception as e:
        print(f"Geometric plot error: {str(e)}")
        return None


def create_data_plot(viz_analysis):
    """Create data visualizations"""
    try:
        data_str = viz_analysis.get('data', '') + " " + viz_analysis.get('expression', '')
        numbers = re.findall(r'-?\d+\.?\d*', data_str)
        values = [float(n) for n in numbers if n]
        
        if not values or len(values) < 2:
            return None
        
        fig = go.Figure()
        fig.add_trace(go.Bar(
            x=list(range(1, len(values) + 1)),
            y=values,
            marker=dict(
                color=values,
                colorscale='Viridis',
                showscale=True,
                line=dict(color='white', width=2)
            ),
            text=[f'{v:.1f}' for v in values],
            textposition='auto',
            name='Values'
        ))
        
        fig.update_layout(
            title='Data Visualization',
            xaxis_title='Index',
            yaxis_title='Value',
            template='plotly_white',
            height=500,
            showlegend=False
        )
        
        return fig.to_html(full_html=False)
        
    except Exception as e:
        print(f"Data plot error: {str(e)}")
        return None


def create_statistical_plot(viz_analysis):
    """Create statistical visualizations"""
    try:
        data_str = viz_analysis.get('data', '') + " " + viz_analysis.get('expression', '')
        numbers = re.findall(r'-?\d+\.?\d*', data_str)
        values = [float(n) for n in numbers if n]
        
        if not values or len(values) < 2:
            return None
        
        fig = go.Figure()
        fig.add_trace(go.Histogram(
            x=values,
            marker=dict(
                color='#4361EE',
                line=dict(color='white', width=1)
            ),
            opacity=0.75,
            name='Distribution'
        ))
        
        mean_val = np.mean(values)
        fig.add_vline(
            x=mean_val,
            line_dash="dash",
            line_color="red",
            line_width=3,
            annotation_text=f"Mean: {mean_val:.2f}",
            annotation_position="top"
        )
        
        fig.update_layout(
            title='Statistical Distribution',
            xaxis_title='Value',
            yaxis_title='Frequency',
            template='plotly_white',
            height=500,
            showlegend=True
        )
        
        return fig.to_html(full_html=False)
        
    except Exception as e:
        print(f"Statistical plot error: {str(e)}")
        return None


def create_vector_plot(viz_analysis):
    """Create vector visualizations"""
    try:
        expression = viz_analysis.get('expression', '') + " " + viz_analysis.get('data', '')
        numbers = re.findall(r'-?\d+\.?\d*', expression)
        
        if len(numbers) >= 2:
            x_comp = float(numbers[0])
            y_comp = float(numbers[1])
            
            fig = go.Figure()
            fig.add_trace(go.Scatter(
                x=[0, x_comp],
                y=[0, y_comp],
                mode='lines+markers',
                line=dict(color='#E63946', width=4),
                marker=dict(size=[10, 15], color='#E63946'),
                name=f'Vector ({x_comp}, {y_comp})'
            ))
            
            fig.add_annotation(
                x=x_comp,
                y=y_comp,
                ax=0,
                ay=0,
                xref='x',
                yref='y',
                axref='x',
                ayref='y',
                showarrow=True,
                arrowhead=2,
                arrowsize=1.5,
                arrowwidth=3,
                arrowcolor='#E63946'
            )
            
            fig.update_xaxes(showgrid=True, gridwidth=1, gridcolor='LightGray', zeroline=True, zerolinewidth=2)
            fig.update_yaxes(showgrid=True, gridwidth=1, gridcolor='LightGray', zeroline=True, zerolinewidth=2)
            
            magnitude = np.sqrt(x_comp**2 + y_comp**2)
            
            fig.update_layout(
                title=f'Vector ({x_comp}, {y_comp}) | Magnitude: {magnitude:.2f}',
                xaxis_title='x',
                yaxis_title='y',
                template='plotly_white',
                height=500,
                yaxis=dict(scaleanchor="x", scaleratio=1),
                showlegend=True
            )
            
            return fig.to_html(full_html=False)
        
        return None
        
    except Exception as e:
        print(f"Vector plot error: {str(e)}")
        return None


@app.route('/upload_and_extract', methods=['POST'])
def upload_and_extract():
    if 'api_key' not in session:
        return jsonify({'error': 'API key missing'})

    file = request.files.get('file')
    if not file:
        return jsonify({'error': 'No file provided'})

    filename = file.filename.lower()

    if filename.endswith(('jpg', 'jpeg', 'png')):
        return jsonify({'error': 'Image AI extraction not supported with Groq.'})
    elif filename.endswith('pdf'):
        return extract_from_pdf(file)
    elif filename.endswith('docx'):
        return extract_from_word(file)
    elif filename.endswith('pptx'):
        return extract_from_ppt(file)
    else:
        return jsonify({'error': 'Unsupported file type. Supported: PDF, DOCX, PPTX'})


def extract_from_pdf(file):
    reader = PdfReader(file)
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return extract_math_from_text(text)


def extract_from_word(file):
    doc = Document(file)
    text = "\n".join([para.text for para in doc.paragraphs])
    return extract_math_from_text(text)


def extract_from_ppt(file):
    presentation = pptx.Presentation(file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return extract_math_from_text(text)


def extract_math_from_text(text):
    client = Groq(api_key=session['api_key'])
    prompt = f"Extract only math questions (no explanations) from this text:\n{text}"

    try:
        completion = client.chat.completions.create(
            model=MODEL_NAME,
            messages=[{"role": "user", "content": prompt}],
            temperature=0
        )
        return jsonify({'text': completion.choices[0].message.content.strip()})
    except Exception as e:
        return jsonify({'error': f'Groq API Error: {str(e)}'})


def extract_json_from_response(response_text):
    """Extract JSON from AI response"""
    json_pattern = r'\{[^{}]*(?:\{[^{}]*\}[^{}]*)*\}'
    match = re.search(json_pattern, response_text, re.DOTALL)
    
    if match:
        try:
            return json.loads(match.group(0))
        except:
            pass
    
    return {'can_visualize': False, 'viz_type': 'none'}


if __name__ == '__main__':
    app.run(debug=True)